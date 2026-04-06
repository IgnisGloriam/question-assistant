import os

def extract_text(file_path: str) -> str:
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Файл не найден: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    parsers = {
        ".txt": _parse_txt,
        ".md": _parse_txt,
        ".docx": _parse_docx,
        ".pdf": _parse_pdf,
        ".pptx": _parse_pptx,
    }

    parser = parsers.get(ext)
    if parser is None:
        supported = ", ".join(parsers.keys())
        raise ValueError(
            f"Неподдерживаемый формат '{ext}'. "
            f"Поддерживаются: {supported}"
        )

    text = parser(file_path)
    text = _clean_text(text)

    return text


def _parse_txt(file_path: str) -> str:
    for encoding in ["utf-8", "cp1251", "latin-1"]:
        try:
            with open(file_path, encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue

    # читаем с игнорированием ошибок, если не вышло нормально
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_docx(file_path: str) -> str:
    from docx import Document
    
    doc = Document(file_path)
    parts = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_text.append(cell_text)
            if row_text:
                parts.append(" | ".join(row_text))

    return "\n".join(parts)


def _parse_pdf(file_path: str) -> str:
    import fitz  # PyMuPDF
    
    doc = fitz.open(file_path)
    parts = []

    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            parts.append(text.strip())

    doc.close()

    return "\n\n".join(parts)


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    
    prs = Presentation(file_path)
    parts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # Текст из таблиц на слайдах
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        slide_texts.append(" | ".join(row_text))

        if slide_texts:
            header = f"[Слайд {slide_num}]"
            parts.append(header + "\n" + "\n".join(slide_texts))

    return "\n\n".join(parts)



def _clean_text(text: str) -> str:
    """Очистка и нормализация текста"""
    
    import re
    text = re.sub(r"[ \t]+", " ", text)

    # Убираем множественные пустые строки (оставляем максимум одну)
    text = re.sub(r"\n{3,}", "\n\n", text)

    # Убираем пробелы в начале и конце строк
    lines = [line.strip() for line in text.split("\n")]
    text = "\n".join(lines)

    # убираем множественные пробелы
    text = re.sub(r'[ \t]+', ' ', text)

    # Убираем пробелы в начале и конце всего текста
    text = text.strip()

    return text





#  тест

if __name__ == "__main__":
    

    path = 'test_text/1 engl.docx'

    try:
        result = extract_text(path)
        print(f"Файл: {path}")
        print(f"Символов: {len(result)}")
        print(f"Слов: {len(result.split())}")
        print("-" * 50)
        
        preview = result[:1000]
        if len(result) > 1000:
            preview += "\n\n... (текст обрезан для превью)"
        print(preview)

    except (FileNotFoundError, ValueError, ImportError) as e:
        print(f"Ошибка: {e}")